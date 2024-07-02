import os
import io
import time
import json
import logging
from io import BytesIO
import base64
import pandas as pd
import streamlit as st
from PIL import Image
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.src.lida import TextGenerationConfig, Manager
from llmx import llm
#from src.lida.components import Manager
from utils.src.lida.datamodel import TextGenerationConfig
from llmx import TextGenerator


logging.basicConfig(level=logging.INFO)

def load_settings(env_path) -> dict:
    load_dotenv(env_path)
    return os.getenv('OPENAI_API_KEY')

class FeatureDescriber:
    def __init__(self) -> None:
        pass

    def describe(
        self,
        summary: dict,
        goal: str,
        base64_image,
        text_gen: TextGenerator,
        textgen_config: TextGenerationConfig,
    ) -> dict:
        """
        Generate descriptions for the given data features based on the summary and goal.

        Args:
            summary (dict): Summary of the dataset in JSON format.
            goal (str): The goal or purpose of the data description.
            text_gen (TextGenerator): Text generator object.
            textgen_config (TextGenerationConfig): Text generation configuration.

        Returns:
            dict: Generated descriptions for the data features.
        """

        system_prompt_base = f"""
        あなたは、経験豊富なデータアナリストです。「{goal}」の目的に沿ったデータを説明してください。
        必ず日本語で出力してください。体言止めにしてください。説明には改行などのフォーマットやラベルを絶対に含まないでください。
        フィールドの情報: {summary['fields']}
        """

        user_prompt = f"""
        あなたは、経験豊富なデータアナリストです。以下の項目について説明してください。
        1. データセットの目的: データセットの目的を説明する文章は末尾に〜データセットとなるようにしてください。要点を簡潔に表現してください。
        2. グラフの解釈: グラフが何を示しているのか、視覚的な観点から説明してください。要点を簡潔に表現してください。
        3. 主要なインサイト: データセットやグラフから得られる具体的な発見やトレンドを詳細に説明してください。具体的な数値や変化に触れて事実を述べてください。
        

        それぞれの項目で説明する内容が重複しないように注意してください。以下の形式に従ってください：
        {{
            "dataset_purpose": "データセットの目的",
            "graph_interpretation": "グラフの解釈",
            "key_insights": "主要なインサイト"
        }}

        以下の画像は作成したグラフを示しています。参考にしてください：
        data:image/png;base64,{base64_image}
        """
        messages = [

            {"role": "system", "content": system_prompt_base.strip()},
            {"role": "user", "content":user_prompt.strip() }   

        ]
        response = text_gen.generate(messages=messages, config=textgen_config)

        print("Response content:", response.text[0]['content'])

        try:
            feature_descriptions = json.loads(response.text[0]['content'])
        except json.JSONDecodeError as e:
            logging.error("JSONDecodeError: %s", e)
            logging.error("Invalid response content: %s", response.text[0]['content'])
            return {}

        return feature_descriptions



class GraphGeneration():
    def __init__(self, openai_key):
        self.manager = Manager(text_gen=llm("openai", api_key=openai_key))
        self.feature_describer = FeatureDescriber()

    def generate_summary(self, manager,uploaded_file_path, summary_method, model, temperature, use_cache):
        start_time = time.time()
        textgen_config = TextGenerationConfig(n=1, temperature=temperature, model=model, use_cache=use_cache)
        summary = manager.summarize(uploaded_file_path, summary_method=summary_method, textgen_config=textgen_config)
        time_taken = time.time() - start_time
        logging.info(f"Response: {summary}")
        logging.info(f"Time Taken: {time_taken}")
        return summary

    def generate_goals(self,manager,summary, num_goals, model, temperature, use_cache):
        start_time = time.time()
        textgen_config = TextGenerationConfig(n=num_goals, temperature=temperature, model=model, use_cache=use_cache)
        goals = manager.goals(summary, n=num_goals, textgen_config=textgen_config)
        time_taken = time.time() - start_time
        logging.info(f"Response: {goals}")
        logging.info(f"Time Taken: {time_taken}")
        return goals

    def generate_visualizations(self,manager,summary, goal, model, num_visualizations, temperature, use_cache, library):
        start_time = time.time()
        textgen_config = TextGenerationConfig(n=num_visualizations, temperature=temperature, model=model, use_cache=use_cache)
        visualizations = manager.visualize(summary=summary, goal=goal, textgen_config=textgen_config, library=library)
        time_taken = time.time() - start_time
        logging.info(f"Time Taken: {time_taken}")
        return visualizations

    def edit_chart(self, manager, summary, model, temperature, use_cache, code, instructions, library):
        textgen_config = TextGenerationConfig(n=1, temperature=temperature, model=model, use_cache=use_cache)
        edited_charts = manager.edit(code=code, summary=summary, instructions=instructions, library=library, textgen_config=textgen_config)
        return edited_charts

    def describe_features(
        self,
        manager, 
        summary: dict,
        goal: str,
        base64_image,
        textgen_config: TextGenerationConfig = TextGenerationConfig(),
    ) -> dict:
        """
        Generate descriptions for the given data features based on the summary and goal.

        Args:
            summary (dict): Summary of the dataset in JSON format.
            goal (str): The goal or purpose of the data description.
            textgen_config (TextGenerationConfig, optional): Text generation configuration. Defaults to TextGenerationConfig().

        Returns:
            dict: Generated descriptions for the data features.
        """
        return self.feature_describer.describe(
            summary=summary,
            goal=goal,
            base64_image=base64_image,
            textgen_config=textgen_config,
            text_gen=manager.text_gen,
        )

class VisualizationProcessor:
    @staticmethod
    def process_summary(summary):
        fields = summary["fields"]
        nfields = []
        for field in fields:
            flatted_fields = {"column": field["column"]}
            for row in field["properties"].keys():
                if row != "samples":
                    flatted_fields[row] = field["properties"][row]
                else:
                    flatted_fields[row] = str(field["properties"][row])
            nfields.append(flatted_fields)
        nfields_df = pd.DataFrame(nfields)
        return nfields_df

    @staticmethod
    def display_visualizations(visualizations, titles):
        for idx, viz in enumerate(visualizations):
            if viz.raster:
                imgdata = base64.b64decode(viz.raster)
                img = Image.open(io.BytesIO(imgdata))
                st.image(img, caption=f'Visualization {idx + 1}', use_column_width=True)
        
        selected_title = st.selectbox('選択して詳細表示', index=0, options=titles, key="selected_viz_title")
        return selected_title

    @staticmethod
    def render_visualizations(visualizations):
        file_paths = []  # リストとして初期化
        for idx, viz in enumerate(visualizations):
            if viz.raster:
                imgdata = base64.b64decode(viz.raster)
                img = Image.open(io.BytesIO(imgdata))
                file_path = f'temp_visualization_{idx + 1}.png'
                img.save(file_path)
                file_paths.append(file_path)
                """
                if os.name == 'posix':  # macOSやLinuxの場合
                    if os.uname().sysname == 'Darwin':  # macOS
                        os.system(f'open {file_path}')
                    else:  # Linux
                        os.system(f'xdg-open {file_path}')
                elif os.name == 'nt':  # Windowsの場合
                    os.startfile(file_path)
                """
        if not file_paths:
            raise ValueError("No visualizations were generated.")
        
        return file_paths


    @staticmethod
    def display_selected_visualization(viz_titles, visualizations, selected_title):
        selected_viz = None
        selected_index = None
        if selected_title is not None:
            selected_index = viz_titles.index(selected_title)
            selected_viz = visualizations[selected_index]
            if selected_viz.raster:
                imgdata = base64.b64decode(selected_viz.raster)
                img = Image.open(io.BytesIO(imgdata))
                st.image(img, caption=selected_title, use_column_width=True)
        return selected_viz, selected_index

class PPTXGenerator:
    @staticmethod
    def save_to_pptx(descriptions, image):
        prs = Presentation("./../utils/template.pptx")
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        content_slide = prs.slide_layouts[6]
        slide = prs.slides.add_slide(content_slide)

        bg = slide.shapes.add_picture('./../template/title_content.png', 0, 0, height=prs.slide_height)
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        tb_left = Inches(0.6)
        tb_top = Inches(1.5)
        tb_width = Inches(5.8)
        tb_height = Inches(7)
        textbox = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
        tf = textbox.text_frame
        tf.text = ""
        tf.auto_size = True
        tf.word_wrap = True

        for points in descriptions.values():
            p = tf.add_paragraph()
            p.text = points
            p.level = 0
            p.space_before = Pt(14)
            p.space_after = Pt(14)
            p.font.size = Pt(16)
            p.font.name = 'Meiryo'
            tf.add_paragraph()

        img_left = Inches(6.5)
        img_top = Inches(1.25)
        img_height = Inches(5)
        if image.raster:
            imgdata = base64.b64decode(image.raster)
            filename = "temp_image.png"
            with open(filename, "wb") as f:
                f.write(imgdata)
            slide.shapes.add_picture(filename, img_left, img_top, height=img_height)
            os.remove(filename)

        binary_file = BytesIO()
        prs.save(binary_file)
        binary_file.seek(0)
        return binary_file


def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")



if __name__ == "__main__":

    ad = os.getcwd()
    openai_key = load_settings('.env')
    uploaded_file_path = os.path.join(ad, 'data', 'Survey_Results.csv')
    temperature = 1.0
    summary_method = 'llm'  
    models = 'gpt-4o' 
    use_cache = False
    num_goals = 3
    num_viz=1
    visualization_libraries = ["seaborn", "matplotlib", "plotly"]

    gg = GraphGeneration(openai_key)
    vp = VisualizationProcessor()

    summary = gg.generate_summary(
        uploaded_file_path=uploaded_file_path,
        summary_method=summary_method,
        model=models,
        temperature=temperature,
        use_cache=use_cache
    )

    goals = gg.generate_goals(
        summary=summary,
        num_goals=num_goals,
        model=models,
        temperature=temperature,
        use_cache=use_cache
    )

    visualization = gg.generate_visualizations(
        summary=summary, 
        goal=goals[0], 
        model=models, 
        num_visualizations=num_viz, 
        temperature=temperature, 
        use_cache=use_cache, 
        library=visualization_libraries[0])

    file_path = vp.render_visualizations(visualization)

    base64_image = encode_image(file_path)
    features = gg.describe_features(
        summary=summary,
        goal=goals[0],
        base64_image=base64_image
    )
    print('==================')
    print('データの説明')
    print(features)
    print('==================')
