from pptx import Presentation
from pptx.util import Inches, Pt
import logging
import os
import time
from io import BytesIO
from pptx.enum.text import PP_ALIGN
from PIL import Image
import base64
import tempfile

def generate_ppt(title: str, content: dict, num_of_slides: int, img_path: list, generated_graphs: list) -> Presentation:
    """タイトルと内容からPPTを生成
    
    Args:
        title (str): タイトル
        content (dict): 内容
        num_of_slides (int): スライド数
        img_path (list): 画像のパス
        
    Returns: 
        Presentation: 生成されたPPT
        
    """ 
    
    start = time.time()
    
    # Create Presentation
    prs = Presentation("app/utils/template.pptx")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    bg = slide.shapes.add_picture('app/template/cover.png', 0, 0, height = prs.slide_height)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    ppt_title = slide.shapes.title
    ppt_title.top = Inches(3)
    ppt_title.left = Inches(0.6)
    ppt_title.width = Inches(12)
    ppt_title.height = Inches(1)
    ppt_title.text = title
    ppt_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    ppt_title.text_frame.paragraphs[0].font.name = 'Meiryo'
    
    # Content Slides
    for i in range(num_of_slides):
        content_slide = prs.slide_layouts[5]
        slide = prs.slides.add_slide(content_slide)
        
        # Add background image
        bg = slide.shapes.add_picture('app/template/title_content.png', 0, 0, height = prs.slide_height)
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
        
        # Insert slide title
        slide_title = slide.shapes.title
        slide_title.text = content[f'slide{i+1}']['title']
        slide_title.top = Inches(0.2)
        slide_title.left = Inches(0.6)
        slide_title.width = Inches(12)
        slide_title.height = Inches(0.5)
        slide_title.text_frame.paragraphs[0].font.name = 'Meiryo'
        slide_title.text_frame.paragraphs[0].font.size = Pt(24)
        slide_title.text_frame.paragraphs[0].font.bold = True


        # Set default layout (text on left, image/graph on right)
        tb_left = Inches(0.6)
        tb_top = Inches(1)
        tb_width = Inches(6)
        tb_height = Inches(5.5)
        
        img_left = Inches(7)
        img_top = Inches(1)
        img_width = Inches(6)
        img_height = Inches(5.5)

        # Check if image exists and add it
        if i < len(img_path) and img_path[i] is not None:
            image = Image.open(img_path[i])
            image_ratio = image.width / image.height
            
            if image_ratio > 2:  # Wide image, place it at the bottom
                img_left = Inches(0.6)
                img_top = Inches(4)
                img_width = Inches(12)
                img_height = Inches(3)
                
                tb_height = Inches(3)  # Adjust text box height
            
            picture = slide.shapes.add_picture(img_path[i], img_left, img_top, width=img_width, height=img_height)
        
        # Check if graph exists and add it
        elif i < len(generated_graphs):
            graph_data = generated_graphs[i]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                if isinstance(graph_data, dict) and 'base64_image' in graph_data:
                    temp_file.write(base64.b64decode(graph_data['base64_image']))
                else:
                    temp_file.write(base64.b64decode(graph_data))
                temp_file_path = temp_file.name

            try:
                slide.shapes.add_picture(temp_file_path, img_left, img_top, width=img_width, height=img_height)
            finally:
                os.unlink(temp_file_path)
            
            
            # Insert content into textbox
            tb = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
            tf = tb.text_frame
            tf.word_wrap = True
            for points in content[f'slide{i+1}']['content']:
                p = tf.add_paragraph()
                p.text = f'・{points}'
                p.font.size = Pt(18)
                p.font.name = 'Meiryo'
                p = tf.add_paragraph()
                p.text = ''
    
  
    # Add final slides
    final_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(final_slide)
    bg = slide.shapes.add_picture('app/template/final_page.png', 0, 0, height = prs.slide_height)
    
    # Save ppt in binary format
    binary_file = BytesIO()
    prs.save(binary_file)
    
    time_taken = time.time() - start
    
    logging.info(f"PPT Generated")
    logging.info(f"Time taken: {time_taken}")
    
    return binary_file

if __name__ == "__main__":
    title = "AIの応用"
    content = {'slide1': {'title': '開発プロセスの概要', 'content': ['文書の解析から開発プロセスがスタートします。', 'プレゼンテーションの構成案の提案が行われます。', 'NLP（自然言語処理）を利用した文章生成が実施されます。', '最終的にページ分割を行い、文書を完成させます。'], 'graphic_prompt': '文書を分析している人、アイデアを練っている人、コンピューターで文章を生成している人、そして文書を整理している人のイラスト'}}
    img_path = ["./tmp/slide1.png"]
    binary_ppt = generate_ppt(title=title, content=content, num_of_slides=1, img_path=img_path)
    # Save binary ppt to pptx
    with open(f"../output/{title}.pptx", "wb") as f:
        f.write(binary_ppt.getbuffer())
    logging.info(f"Saved ppt")
